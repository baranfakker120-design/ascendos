#!/usr/bin/env python3
"""One-shot: compose AscendOS_Presentation.pdf + .pptx from framed board PNGs."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

ROOT = Path(__file__).resolve().parents[1]
BOARDS = ROOT / "presentation" / "boards"
OUT_DIR = ROOT / "presentation" / "final"
ARTIFACTS = Path("/opt/cursor/artifacts/presentation")

# Widescreen 16:9
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
SLIDE_W_PT, SLIDE_H_PT = SLIDE_W_IN * 72, SLIDE_H_IN * 72

SLIDES: list[tuple[str, str, str]] = [
    ("01-hero-iphone-today.png", "Today", "Your operating system for the workday"),
    ("01b-hero-login-real.png", "Welcome", "Quiet confidence. Premium by default."),
    ("02-dual-today-coach.png", "Today × Coach", "Plan the day. Lead the conversation."),
    ("02b-dual-team-contacts.png", "Structure × Relationships", "See the tree. Move the pipeline."),
    ("03-triple-workflow.png", "Lead → Plan → Coach", "One operating loop for recruiting days"),
    ("04-feature-structure-tree.png", "Structure Tree", "Living genealogy with ranks and leadership signals"),
    ("04b-feature-coach.png", "Ascent Coach", "Context-first mentoring. One clear next action."),
    ("04c-feature-qualifications.png", "Qualifications", "AP, ranks, and Team Leader progress"),
    ("05-appstore-today.png", "App Store · Today", "Minimal. Focused. Ready to ship."),
    ("05b-appstore-profile.png", "App Store · Profile", "Identity, rank, and recognition"),
    ("05c-appstore-login-real.png", "App Store · Login", "Real production login surface"),
    ("06-dark-premium-coach.png", "Dark Premium · Coach", "Cinematic product presence"),
    ("06b-dark-premium-team.png", "Dark Premium · Team", "Leadership at a glance"),
    ("07-investor-grid.png", "Product Surface", "Investor overview — nine key screens"),
]


def fit_image(path: Path, max_w: int, max_h: int) -> Image.Image:
    img = Image.open(path).convert("RGB")
    img.thumbnail((max_w, max_h), Image.Resampling.LANCZOS)
    return img


def build_pdf(out: Path) -> None:
    c = canvas.Canvas(str(out), pagesize=(SLIDE_W_PT, SLIDE_H_PT))

    # Title slide
    c.setFillColorRGB(0.06, 0.07, 0.08)
    c.rect(0, 0, SLIDE_W_PT, SLIDE_H_PT, fill=1, stroke=0)
    c.setFillColorRGB(0.72, 0.58, 0.35)
    c.setFont("Helvetica-Bold", 18)
    c.drawString(72, SLIDE_H_PT - 100, "ASCENDOS")
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 54)
    c.drawString(72, SLIDE_H_PT - 180, "Build a better tomorrow.")
    c.setFillColorRGB(0.75, 0.75, 0.76)
    c.setFont("Helvetica", 22)
    c.drawString(72, SLIDE_H_PT - 230, "Product presentation · iPhone 16 Pro renders")
    c.setFont("Helvetica", 14)
    c.drawString(72, 72, "Confidential · For investors, partners & events")
    c.showPage()

    for filename, title, subtitle in SLIDES:
        path = BOARDS / filename
        if not path.exists():
            continue
        c.setFillColorRGB(0.04, 0.045, 0.05)
        c.rect(0, 0, SLIDE_W_PT, SLIDE_H_PT, fill=1, stroke=0)

        # Header
        c.setFillColorRGB(0.72, 0.58, 0.35)
        c.setFont("Helvetica-Bold", 11)
        c.drawString(48, SLIDE_H_PT - 36, "ASCENDOS")
        c.setFillColorRGB(1, 1, 1)
        c.setFont("Helvetica-Bold", 26)
        c.drawString(48, SLIDE_H_PT - 68, title)
        c.setFillColorRGB(0.7, 0.7, 0.72)
        c.setFont("Helvetica", 13)
        c.drawString(48, SLIDE_H_PT - 90, subtitle)

        max_w = int(SLIDE_W_PT - 96)
        max_h = int(SLIDE_H_PT - 140)
        img = fit_image(path, max_w, max_h)
        x = (SLIDE_W_PT - img.width) / 2
        y = 36 + (max_h - img.height) / 2
        c.drawImage(ImageReader(img), x, y, width=img.width, height=img.height, mask="auto")
        c.showPage()

    # Closing
    c.setFillColorRGB(0.06, 0.07, 0.08)
    c.rect(0, 0, SLIDE_W_PT, SLIDE_H_PT, fill=1, stroke=0)
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 48)
    c.drawCentredString(SLIDE_W_PT / 2, SLIDE_H_PT / 2 + 20, "AscendOS")
    c.setFillColorRGB(0.72, 0.58, 0.35)
    c.setFont("Helvetica", 20)
    c.drawCentredString(SLIDE_W_PT / 2, SLIDE_H_PT / 2 - 24, "Build a better tomorrow.")
    c.showPage()
    c.save()


def build_pptx(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    def add_bg(slide, r, g, b):
        shape = slide.shapes.add_shape(
            1,  # rectangle
            Inches(0),
            Inches(0),
            Inches(SLIDE_W_IN),
            Inches(SLIDE_H_IN),
        )
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        # send to back
        spTree = slide.shapes._spTree
        sp = shape._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def add_text(slide, left, top, width, height, text, size, bold=False, color=(255, 255, 255), align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        run.font.name = "Helvetica"
        return box

    # Title
    s = prs.slides.add_slide(blank)
    add_bg(s, 15, 18, 20)
    add_text(s, Inches(1), Inches(1.2), Inches(10), Inches(0.4), "ASCENDOS", 18, True, (184, 147, 90))
    add_text(s, Inches(1), Inches(2.0), Inches(11), Inches(1), "Build a better tomorrow.", 48, True)
    add_text(
        s,
        Inches(1),
        Inches(3.1),
        Inches(11),
        Inches(0.5),
        "Product presentation · iPhone 16 Pro renders",
        20,
        False,
        (190, 190, 192),
    )
    add_text(
        s,
        Inches(1),
        Inches(6.6),
        Inches(11),
        Inches(0.4),
        "Confidential · For investors, partners & events",
        12,
        False,
        (140, 140, 144),
    )

    for filename, title, subtitle in SLIDES:
        path = BOARDS / filename
        if not path.exists():
            continue
        s = prs.slides.add_slide(blank)
        add_bg(s, 10, 11, 13)
        add_text(s, Inches(0.7), Inches(0.25), Inches(4), Inches(0.3), "ASCENDOS", 12, True, (184, 147, 90))
        add_text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.45), title, 26, True)
        add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.35), subtitle, 14, False, (180, 180, 184))

        # Fit image into remaining area
        img = Image.open(path)
        iw, ih = img.size
        max_w, max_h = 12.0, 5.9
        scale = min(max_w / (iw / 96), max_h / (ih / 96), max_w / (iw / 72) * 0.01 + 1)
        # Use inches from pixel aspect
        aspect = iw / ih
        if max_w / max_h > aspect:
            disp_h = max_h
            disp_w = disp_h * aspect
        else:
            disp_w = max_w
            disp_h = disp_w / aspect
        left = (SLIDE_W_IN - disp_w) / 2
        top = 1.35 + (max_h - disp_h) / 2
        s.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(disp_w), Inches(disp_h))

    # Close
    s = prs.slides.add_slide(blank)
    add_bg(s, 15, 18, 20)
    add_text(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1), "AscendOS", 48, True, (255, 255, 255), PP_ALIGN.CENTER)
    add_text(
        s,
        Inches(1),
        Inches(3.9),
        Inches(11.3),
        Inches(0.5),
        "Build a better tomorrow.",
        20,
        False,
        (184, 147, 90),
        PP_ALIGN.CENTER,
    )

    prs.save(str(out))


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ARTIFACTS.mkdir(parents=True, exist_ok=True)

    pdf_path = OUT_DIR / "AscendOS_Presentation.pdf"
    pptx_path = OUT_DIR / "AscendOS_Presentation.pptx"
    build_pdf(pdf_path)
    build_pptx(pptx_path)

    # Also place copies at presentation root and artifacts for easy pickup
    for dest_dir in (ROOT / "presentation", ARTIFACTS):
        dest_dir.mkdir(parents=True, exist_ok=True)
        for src in (pdf_path, pptx_path):
            target = dest_dir / src.name
            target.write_bytes(src.read_bytes())

    print(pdf_path.resolve())
    print(pptx_path.resolve())
    print((ROOT / "presentation" / "AscendOS_Presentation.pdf").resolve())
    print((ROOT / "presentation" / "AscendOS_Presentation.pptx").resolve())
    print((ARTIFACTS / "AscendOS_Presentation.pdf").resolve())
    print((ARTIFACTS / "AscendOS_Presentation.pptx").resolve())


if __name__ == "__main__":
    main()
